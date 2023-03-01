import getpass
import sqlite3
import pptx

import pandas as pd
import datetime as dt

from pptx.util import Pt
from MailReader import get_programacao
from frequencia import get_frequencia


print("---------- Hello ----------")


def verificar_falta(df_ativos):
    """
    Remove ativos de acordo com abstinência e regras de escala dinâmica
    :df_ativos:
    :return:
    """



    return df_ativos


def create_slide(presentation, slide_atual, title, id_tamplate=0):
    print('Criando novo slide')
    slide_template = presentation.slides[id_tamplate]

    new_slide = presentation.slides.add_slide(slide_template.slide_layout)
    slide_atual += 1
    for shape in slide_template.shapes:
        if shape.has_table:
            new_shape = new_slide.shapes.add_table(
                11,
                1,
                shape.left,
                shape.top,
                shape.width,
                shape.height
            )

            table = shape.table
            new_table = new_shape.table

            tbl = new_shape._element.graphic.graphicData.tbl
            tbl[0][-1].text = '{6E25E649-3F16-4E02-A733-19D2CDBF48F0}'

            print('Alterar fonte')
            for row in new_table.rows:
                for cell in row.cells:
                    if not cell.text_frame:
                        cell.text_frame = cell._element.add_now_txBody().add_new_p()

                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)

            new_table.first_row = table.first_row
            new_table.first_col = table.first_col
            new_table.last_row = table.last_row
            new_table.last_col = table.last_col
        elif shape.has_text_frame:
            new_shape = new_slide.shapes.add_textbox(
                shape.left,
                shape.top,
                shape.width,
                shape.height
            )
            new_shape.text = title

    return slide_atual


# obtendo o nome do usuário atual
if __name__ == '__main__':
    user = getpass.getuser()
    # Definindo o diretório para as planilhas
    path = f'C:\\Users\\{user}\\Procter and Gamble\\Grupo Check List - Bases de dados\\'

    # Definindo diretório para planilhas de controle (teste power apps)
    # C:\Users\wanderson.wf\Procter and Gamble\Gutemberg, Borges - Projeto CRUD - Teste
    path_controle = f'C:\\Users\\{user}\\Procter and Gamble\\Gutemberg, Borges - Projeto CRUD - Teste\\'

    # Definindo acessos para a planilha de cadastros
    cadastro_file = 'planilhas\\df_Cadastro.xlsx'

    # Tabela com o cadastro de códigos produtos e materia-prima
    df_cadastro_codigo = pd.read_excel(path + cadastro_file, sheet_name='dCódigo', dtype=str)

    # Definincdo acesso para tabela de feriados nacionais
    feriados_file = 'planilhas\\df_Feriado.xlsx'

    # Tabela com os fériados nacionais relevantes
    df_feriados = pd.read_excel(path + feriados_file)

    # Definindo acesso para tabela de Ativos
    ativos_file = 'teste_df_Ativos_Yusen.xlsx'
    # Obtendo dados relacionados aos Ativos
    df_ativos_yusen = pd.read_excel(path_controle + ativos_file, dtype=str)

    # ## 'df_programacao' consolidado.
    df_programacao, _ = get_programacao(df_cadastro_codigo)

    df_programacao['Prioridade'] = ''
    df_programacao.loc[
        df_programacao['Operação'] == 'SHERINKADORA',
        'Prioridade'
    ] = 'X'

    df_programacao.loc[
        df_programacao['Operação'] == 'SERIES I',
        'Prioridade'
    ] = 'W'

    df_programacao.sort_values(by=['Prioridade', 'Operação'],
                               ascending=False, inplace=True)

    # 'df_frequencia' para lógica de frequencia
    df_frequencia = get_frequencia()

    df_prog_today = pd.DataFrame(columns=['Descrição de linha', 'Código', 'Turno', 'Matrícula', 'Data'])

    list_turno = [
        '1º Turno',
        '2º Turno',
        '3º Turno'
    ]

    for turno in list_turno:
        print(turno)
        df_temp_fturno = df_ativos_yusen.loc[df_ativos_yusen['Turno'] == turno]
        df_temp_pturno = df_programacao.loc[df_programacao[turno].notnull()]
        list_geral_func = []

        for indice in df_temp_pturno.index:
            desc_linha, codigo = df_temp_pturno.loc[indice, ['Descrição de linha', 'Código']].values

            df_temp_ffreq = df_temp_fturno.merge(
                df_frequencia.loc[
                    (df_frequencia['Descrição de linha'] == desc_linha) &
                    (df_frequencia['Código'] == codigo),
                    ['Matrícula', 'Frequencia', 'Data']
                ],
                on=['Matrícula'],
                how='left'
            )

            # Lógica de frequência
            df_temp_ffreq.sort_values(['Data', 'Frequencia'], inplace=True, na_position='first')

            operacao, mdo = df_temp_pturno.loc[indice, ['Operação', 'MDO']].values
            if df_temp_pturno.loc[indice, 'Prioridade'] == 'X':
                list_func = list(df_temp_ffreq.loc[
                                     df_temp_ffreq['Habilidade'] == 'Nível 3',
                                     'Matrícula'
                                 ])
            elif df_temp_pturno.loc[indice, 'Prioridade'] == 'W':
                list_func = list(df_temp_ffreq.loc[
                                     (df_temp_ffreq['Habilidade'] == 'Nível 3') |
                                     (df_temp_ffreq['Habilidade'] == 'Nível 2'),
                                     'Matrícula'
                                 ])
            else:
                list_func = list(df_temp_ffreq['Matrícula'])

            repetidos = [func for func in list_func if func in list_geral_func]
            if repetidos:
                for func in repetidos:
                    list_func.remove(func)

            # Lógica de frequencia (no corte)
            if int(mdo) <= len(list_func):
                list_func = list_func[0:int(mdo)]

            list_geral_func += list_func

            for matricula in list_func:
                df_prog_today.loc[len(df_prog_today)] = [
                    desc_linha,
                    codigo,
                    turno,
                    matricula,
                    # dt.date.today().strftime('%d/%m/%Y')
                    dt.date(day=5, month=3, year=2023).strftime('%d/%m/%Y')
                ]

    # Definindo acessi para DB de programação de aeroporto
    db_airport_file = 'bd_sqlite\\db_airport.db'
    conn = sqlite3.connect(path + db_airport_file)

    # Realizando gravação no historico de programação do aeroporto
    # df_prog_today.to_sql('hist_programacao', conn, if_exists='append', index=False)

    # Tabela de histórivo de programaçao (Verificação de gravação)
    df_programacao_historico = pd.read_sql('SELECT * FROM hist_programacao', conn)

    # Fechando o conector
    conn.close()
    print("Programação gerada com sucesso!")

    indice = 0
    ppt = pptx.Presentation('template.pptx')

    for turno in list_turno:
        # indice de controle para slide e tabela usados
        titulo = f'Programação do {turno} na data {df_prog_today["Data"][0]}'
        indice = create_slide(ppt, indice, titulo)
        shape_id = 0

        # informações por turno
        df_temp_pturno = df_prog_today.loc[df_prog_today['Turno'] == turno]
        df_temp_fturno = df_ativos_yusen.loc[df_ativos_yusen['Turno'] == turno]
        df_list_codigos = df_temp_pturno.drop_duplicates([
            'Código', 'Descrição de linha'])[['Código', 'Descrição de linha']]

        # Execução de saida no ppt
        for codigo, desc_linha in df_list_codigos.values:
            # informações por linha/código
            filtered = df_temp_pturno.loc[
                (df_temp_pturno['Código'] == codigo) &
                (df_temp_pturno['Descrição de linha'] == desc_linha)
            ]

            print(' -- ', shape_id)
            if shape_id >= 8:
                indice = create_slide(ppt, indice, titulo)
                shape_id = 0
                print(shape_id)

            fontsize = Pt(12)
            table_ppt = ppt.slides[indice].shapes[shape_id].table
            table_ppt.cell(0, 0).text = f'{codigo}\n{desc_linha}'
            table_ppt.cell(0, 0).text_frame.paragraphs[0].font.size = fontsize
            table_ppt.cell(0, 0).text_frame.paragraphs[1].font.size = fontsize

            linha = 1
            for matricula in filtered['Matrícula'].values:
                table_ppt.cell(linha, 0).text = df_temp_fturno.loc[df_temp_fturno['Matrícula'] == matricula,
                                                                   'Nome do Funcionário'].values[0]
                table_ppt.cell(linha, 0).text_frame.paragraphs[0].font.size = fontsize
                linha += 1
            shape_id += 1

            print(codigo, desc_linha)

    ppt.save('test_1.pptx')
    print("---------------------")
